from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

import fitz
import pytesseract
from docx import Document
from PIL import Image
from pptx import Presentation


@dataclass(frozen=True)
class ExtractedSection:
    marker: str
    text: str
    page_number: int | None = None


def extract_text(path: Path, extension: str) -> list[ExtractedSection]:
    match extension.lower():
        case ".pdf":
            return _extract_pdf(path)
        case ".pptx":
            return _extract_pptx(path)
        case ".docx":
            return _extract_docx(path)
        case ".txt" | ".md" | ".vtt" | ".srt":
            return _extract_plain_text(path)
        case _:
            raise ValueError(f"unsupported_file_type:{extension}")


def _extract_pdf(path: Path) -> list[ExtractedSection]:
    sections: list[ExtractedSection] = []
    with fitz.open(path) as doc:
        for index, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if len(text) < 100:
                text = _ocr_pdf_page(page).strip()
            sections.append(
                ExtractedSection(
                    marker=f"[Page {index}]",
                    page_number=index,
                    text=f"[Page {index}]\n{text}",
                )
            )
    return sections


def _ocr_pdf_page(page: fitz.Page) -> str:
    pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
    image = Image.open(BytesIO(pixmap.tobytes("png")))
    return pytesseract.image_to_string(image)


def _extract_pptx(path: Path) -> list[ExtractedSection]:
    presentation = Presentation(path)
    sections: list[ExtractedSection] = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = list(_slide_shape_text(slide.shapes))
        body = "\n".join(part for part in texts if part.strip()).strip()
        sections.append(
            ExtractedSection(
                marker=f"[Slide {index}]",
                page_number=index,
                text=f"[Slide {index}]\n{body}",
            )
        )
    return sections


def _slide_shape_text(shapes) -> list[str]:
    parts: list[str] = []
    for shape in shapes:
        if hasattr(shape, "shapes"):
            parts.extend(_slide_shape_text(shape.shapes))
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                line = "".join(run.text for run in paragraph.runs).strip()
                if line:
                    parts.append(line)
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                parts.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
    return parts


def _extract_docx(path: Path) -> list[ExtractedSection]:
    document = Document(path)
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            line = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if line:
                paragraphs.append(line)
    return [ExtractedSection(marker="[Document]", text="[Document]\n" + "\n".join(paragraphs))]


def _extract_plain_text(path: Path) -> list[ExtractedSection]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return [ExtractedSection(marker="[Transcript]", text="[Transcript]\n" + text)]
